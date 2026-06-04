#!/usr/bin/env python3
"""Render a Tufte decision matrix from a build-options decision-data.json.

Usage: build_matrix.py <decision-data.json> [--out DIR] [--pdf]

Weighted totals + tie analysis are COMPUTED here via weighting.py (the single source of truth) from the
raw per-judge means in decision-data.json — never read from the file — so the viz always matches the math.
HTML always; PPTX if python-pptx importable; PDF only with --pdf + chromium.
"""
import json, re, sys, pathlib, subprocess, shutil

HERE = pathlib.Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))
import weighting  # noqa: E402

TEMPLATE = HERE / "matrix.template.html"


def slugify(t):
    s = re.sub(r"[^a-zA-Z0-9]+", "-", str(t)).strip("-").lower()
    return s or "decision"


def enrich(data):
    """Compute weightedTotal per option, ranking, and tie analysis from raw score means."""
    criteria = data["criteria"]
    mean_scores = {}
    for o in data["options"]:
        means = {k: v["mean"] for k, v in o.get("scores", {}).items()}
        mean_scores[o["id"]] = means
        o["weightedTotal"] = weighting.weighted_total(means, criteria)
    ranked = sorted(((o["id"], o["weightedTotal"]) for o in data["options"]), key=lambda x: -x[1])
    data["ranking"] = [oid for oid, _ in ranked]
    data["tie"] = weighting.detect_tie(ranked, criteria=criteria, option_scores=mean_scores)
    ok, total = weighting.validate_weights(criteria)
    data["weightsOk"] = ok
    data["weightsTotal"] = total
    return data


def main():
    argv = sys.argv[1:]
    do_pdf = "--pdf" in argv
    argv = [a for a in argv if a != "--pdf"]
    out_dir = pathlib.Path(".")
    if "--out" in argv:
        i = argv.index("--out"); out_dir = pathlib.Path(argv[i + 1]); del argv[i:i + 2]
    data_path = pathlib.Path(argv[0]) if argv else pathlib.Path("decision-data.json")
    out_dir.mkdir(parents=True, exist_ok=True)

    data = enrich(json.loads(data_path.read_text()))
    slug = slugify(data.get("meta", {}).get("title", "decision"))

    payload = json.dumps(data, ensure_ascii=False).replace("</", "<\\/")
    html = TEMPLATE.read_text().replace("/*__DECISION_DATA__*/ null", payload)
    html_path = out_dir / (slug + ".html")
    html_path.write_text(html)
    assert "__DECISION_DATA__" not in html, "placeholder not injected"
    print("HTML", html_path, len(html), "chars")

    try:
        build_pptx(data, out_dir / (slug + ".pptx"))
        print("PPTX", out_dir / (slug + ".pptx"))
    except ImportError:
        print("PPTX_SKIPPED (python-pptx not installed)")

    if do_pdf:
        chrome = (shutil.which("chromium") or shutil.which("chromium-browser")
                  or shutil.which("google-chrome") or shutil.which("chrome"))
        if chrome:
            pdf_path = out_dir / (slug + ".pdf")
            subprocess.run([chrome, "--headless", "--disable-gpu", "--no-sandbox", "--no-pdf-header-footer",
                            "--virtual-time-budget=10000", "--run-all-compositor-stages-before-draw",
                            "--print-to-pdf=" + str(pdf_path), html_path.resolve().as_uri()],
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=120)
            print("PDF", pdf_path if pdf_path.exists() else "PDF_FAILED")
        else:
            print("PDF_SKIPPED (no chromium found)")


def build_pptx(data, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    INK = RGBColor(0x1A, 0x1A, 0x1A); GRN = RGBColor(0x2F, 0x6B, 0x4F); PAPER = RGBColor(0xFF, 0xFF, 0xF8)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF); SOFT = RGBColor(0x55, 0x51, 0x4A)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    by_id = {o["id"]: o for o in data["options"]}
    rec = data.get("recommendation", {}); win = by_id.get(rec.get("optionId"), {})
    # title + recommendation
    s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = PAPER
    t = s.shapes.add_textbox(Inches(.6), Inches(.6), Inches(12), Inches(1.2)).text_frame; t.word_wrap = True
    t.paragraphs[0].text = data["meta"]["title"]
    for r in t.paragraphs[0].runs: r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = INK
    b = s.shapes.add_textbox(Inches(.6), Inches(1.9), Inches(12), Inches(4.8)).text_frame; b.word_wrap = True
    p0 = b.paragraphs[0]; p0.text = "Recommendation: " + win.get("name", rec.get("optionId", "?")) + "  (" + rec.get("confidence", "") + ")"
    for r in p0.runs: r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = GRN
    p1 = b.add_paragraph(); p1.text = rec.get("rationale", "")
    for r in p1.runs: r.font.size = Pt(12); r.font.color.rgb = INK
    # matrix slide
    s2 = prs.slides.add_slide(blank); s2.background.fill.solid(); s2.background.fill.fore_color.rgb = PAPER
    h = s2.shapes.add_textbox(Inches(.5), Inches(.3), Inches(12), Inches(.5)).text_frame
    h.paragraphs[0].text = "Scored options (weighted 0-5)"
    for r in h.paragraphs[0].runs: r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GRN
    opts = [by_id[i] for i in data["ranking"]]
    tbl = s2.shapes.add_table(len(opts) + 1, 4, Inches(.5), Inches(.9), Inches(12.3), Inches(5.5)).table
    for j, htxt in enumerate(["Option", "Lens", "Weighted total", "Adversarial"]):
        c = tbl.cell(0, j); c.text = htxt; c.fill.solid(); c.fill.fore_color.rgb = INK
        rr = c.text_frame.paragraphs[0].runs[0]; rr.font.size = Pt(11); rr.font.bold = True; rr.font.color.rgb = WHITE
    for i, o in enumerate(opts, start=1):
        vals = [o["name"], o.get("lens", ""), "%.2f" % o["weightedTotal"], o.get("adversarial", {}).get("verdict", "")]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = PAPER
            pr = c.text_frame.paragraphs[0]; pr.text = v
            for r in pr.runs: r.font.size = Pt(10); r.font.color.rgb = INK; r.font.bold = (j == 0)
    prs.save(str(out_path))


if __name__ == "__main__":
    main()
