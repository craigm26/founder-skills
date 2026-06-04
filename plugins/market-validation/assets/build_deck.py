#!/usr/bin/env python3
"""Generate a Tufte market-evidence deck from a canonical deck-data.json.

Usage: build_deck.py [deck-data.json] [--out DIR] [--pdf]

- HTML is always produced (self-contained; in-page Save-PDF / Export-PPTX buttons work in a browser).
- PPTX is produced if python-pptx is importable, else PPTX_SKIPPED.
- PDF is produced only with --pdf and a chromium binary, else PDF_SKIPPED.

Output files are named from a slug of meta.title.
"""
import json, re, sys, pathlib, subprocess, shutil

HERE = pathlib.Path(__file__).resolve().parent
TEMPLATE = HERE / "deck.template.html"


def slugify(t):
    s = re.sub(r"[^a-zA-Z0-9]+", "-", str(t)).strip("-").lower()
    return s or "market-evidence"


def strip(t):
    """Drop [n] citations and normalize unicode punctuation for PPTX text."""
    t = re.sub(r"\[(\d+)\]", "", str(t))
    return (t.replace("“", '"').replace("”", '"').replace("’", "'")
            .replace("–", "-").replace("—", "-").replace("∝", "~").strip())


def main():
    argv = sys.argv[1:]
    do_pdf = "--pdf" in argv
    argv = [a for a in argv if a != "--pdf"]
    out_dir = pathlib.Path(".")
    if "--out" in argv:
        i = argv.index("--out"); out_dir = pathlib.Path(argv[i + 1]); del argv[i:i + 2]
    data_path = pathlib.Path(argv[0]) if argv else pathlib.Path("deck-data.json")
    out_dir.mkdir(parents=True, exist_ok=True)

    DATA = json.loads(data_path.read_text())
    slug = slugify(DATA.get("meta", {}).get("title", "market-evidence"))

    # 1. HTML (always)
    payload = json.dumps(DATA, ensure_ascii=False).replace("</", "<\\/")
    html = TEMPLATE.read_text().replace("/*__DECK_DATA__*/ null", payload)
    html_path = out_dir / (slug + ".html")
    html_path.write_text(html)
    assert "__DECK_DATA__" not in html, "data placeholder was not injected"
    print("HTML", html_path, len(html), "chars")

    # 2. PPTX (best-effort)
    pptx_path = out_dir / (slug + ".pptx")
    try:
        build_pptx(DATA, pptx_path)
        print("PPTX", pptx_path)
    except ImportError:
        print("PPTX_SKIPPED (python-pptx not installed)")

    # 3. PDF (best-effort, --pdf)
    if do_pdf:
        chrome = (shutil.which("chromium") or shutil.which("chromium-browser")
                  or shutil.which("google-chrome") or shutil.which("chrome"))
        if chrome:
            pdf_path = out_dir / (slug + ".pdf")
            subprocess.run([chrome, "--headless", "--disable-gpu", "--no-sandbox",
                            "--no-pdf-header-footer", "--virtual-time-budget=12000",
                            "--run-all-compositor-stages-before-draw",
                            "--print-to-pdf=" + str(pdf_path), html_path.resolve().as_uri()],
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=150)
            print("PDF", pdf_path if pdf_path.exists() else "PDF_FAILED")
        else:
            print("PDF_SKIPPED (no chromium found)")


def build_pptx(DATA, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR

    INK = RGBColor(0x1A, 0x1A, 0x1A); GRN = RGBColor(0x2F, 0x6B, 0x4F)
    SOFT = RGBColor(0x55, 0x51, 0x4A); PAPER = RGBColor(0xFF, 0xFF, 0xF8)
    ALERT = RGBColor(0xA2, 0x3B, 0x2E); AMBER = RGBColor(0x9A, 0x63, 0x12)
    BODY = RGBColor(0x2A, 0x27, 0x22)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(BLANK)
        bg = s.background.fill; bg.solid(); bg.fore_color.rgb = PAPER
        return s

    def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        return tf

    def setp(p, text, size, color=BODY, bold=False, italic=False, font="Arial", spacing=1.0):
        p.text = text
        for r in p.runs:
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.italic = italic; r.font.name = font
        p.line_spacing = spacing
        return p

    def heading(s, txt, color=GRN):
        setp(tb(s, .6, .35, 12.1, .55).paragraphs[0], txt, 15, color, bold=True)

    meta = DATA.get("meta", {})
    # title
    s = slide()
    setp(tb(s, .6, 2.1, 12.1, 1.3).paragraphs[0], strip(meta.get("title", "Market Evidence")), 40, INK, bold=True)
    setp(tb(s, .6, 3.5, 12.1, .9).paragraphs[0], strip(meta.get("subtitle", "")), 20, SOFT)
    setp(tb(s, .6, 6.5, 12.1, .6).paragraphs[0],
         meta.get("date", "") + "   ·   " + strip(meta.get("provenance", "")), 11, SOFT)

    # verdict
    v = DATA.get("verdict", {})
    s = slide(); heading(s, "1 · Verdict")
    setp(tb(s, .6, 1.0, 12.2, 1.7).paragraphs[0], strip(v.get("headline", "")), 19, INK, bold=True, spacing=1.05)
    tf = tb(s, .6, 2.75, 8.2, 4.3)
    for i, para in enumerate(v.get("paras", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        setp(p, "•  " + strip(para), 12, BODY, spacing=1.05); p.space_after = Pt(6)
    conf = v.get("confidence", [])
    if conf:
        tbl = s.shapes.add_table(len(conf) + 1, 2, Inches(9.0), Inches(2.75), Inches(3.8), Inches(2.4)).table
        tbl.columns[0].width = Inches(2.7); tbl.columns[1].width = Inches(1.1)
        for j, htxt in enumerate(["Claim", "Conf."]):
            c = tbl.cell(0, j); c.text = htxt; c.fill.solid(); c.fill.fore_color.rgb = INK
            pr = c.text_frame.paragraphs[0].runs[0]; pr.font.size = Pt(10); pr.font.bold = True
            pr.font.color.rgb = WHITE; pr.font.name = "Arial"
        lvlcol = {"HIGH": GRN, "MODERATE": AMBER, "LOW": ALERT}
        for i, cf in enumerate(conf, start=1):
            c0 = tbl.cell(i, 0); c0.fill.solid(); c0.fill.fore_color.rgb = PAPER
            setp(c0.text_frame.paragraphs[0], strip(cf.get("label", "")), 9.5, BODY)
            c1 = tbl.cell(i, 1); c1.fill.solid(); c1.fill.fore_color.rgb = PAPER
            setp(c1.text_frame.paragraphs[0], cf.get("level", ""), 9.5, lvlcol.get(cf.get("level"), BODY), bold=True)

    # pain
    pain = DATA.get("pain", {})
    s = slide(); heading(s, "2 · Pain-point demand")
    tf = tb(s, .6, 1.0, 12.2, 6.2)
    for i, it in enumerate(pain.get("items", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = strip(it.get("point", "")) + "  "
        if p.runs:
            p.runs[0].font.bold = True; p.runs[0].font.color.rgb = INK
        for chunk, kw in [('"' + strip(it.get("quote", "")) + '"  ', dict(italic=True, color=BODY)),
                          ("— " + strip(it.get("who", "")), dict(color=SOFT))]:
            r = p.add_run(); r.text = chunk; r.font.italic = kw.get("italic", False)
            r.font.color.rgb = kw["color"]
        for r in p.runs:
            r.font.size = Pt(11.5); r.font.name = "Arial"
        p.space_after = Pt(7); p.line_spacing = 1.03

    # competitors
    comp = DATA.get("competitors", []); cats = {c["key"]: c for c in DATA.get("categories", [])}
    s = slide(); heading(s, "3 · Competitor validation")
    if comp:
        tbl = s.shapes.add_table(len(comp) + 1, 4, Inches(.4), Inches(.95), Inches(12.5), Inches(6.3)).table
        for col, w in zip(tbl.columns, [1.7, 1.4, 5.2, 4.2]):
            col.width = Inches(w)
        for j, htxt in enumerate(["Company", "Tier", "What it does", "Funding / traction"]):
            c = tbl.cell(0, j); c.text = htxt; c.fill.solid(); c.fill.fore_color.rgb = INK
            rr = c.text_frame.paragraphs[0].runs[0]; rr.font.size = Pt(8.5); rr.font.bold = True
            rr.font.color.rgb = WHITE; rr.font.name = "Arial"
        for i, cp in enumerate(comp, start=1):
            cat = cats.get(cp.get("category")); tier = cat["label"].split(" ")[0] if cat else cp.get("category", "")
            vals = [strip(cp.get("name", "")), tier, strip(cp.get("what", "")), strip(cp.get("funding", ""))]
            for j, val in enumerate(vals):
                c = tbl.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = PAPER
                c.margin_top = Pt(1); c.margin_bottom = Pt(1)
                setp(c.text_frame.paragraphs[0], val, 7.5, BODY, bold=(j == 0))

    # wtp
    wtp = DATA.get("wtp", {})
    s = slide(); heading(s, "4 · Willingness to pay")
    tf = tb(s, .6, 1.0, 12.2, 4.7)
    for i, it in enumerate(wtp.get("items", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  "
        r = p.add_run(); r.text = strip(it.get("point", "")) + ". "; r.font.bold = True; r.font.color.rgb = INK
        r2 = p.add_run(); r2.text = strip(it.get("detail", "")); r2.font.color.rgb = BODY
        for r in p.runs:
            r.font.size = Pt(11.5); r.font.name = "Arial"
        p.space_after = Pt(5); p.line_spacing = 1.03
    if wtp.get("honest"):
        setp(tb(s, .6, 6.0, 12.2, 1.2).paragraphs[0], "Honest read: " + strip(wtp["honest"]), 11, AMBER, italic=True)

    # global
    s = slide(); heading(s, "5 · Global signals")
    tf = tb(s, .6, 1.0, 12.2, 6.2)
    for i, g in enumerate(DATA.get("global", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = g.get("region", "") + ": "; r.font.bold = True; r.font.color.rgb = INK
        r2 = p.add_run(); r2.text = strip(g.get("text", "")); r2.font.color.rgb = BODY
        for r in p.runs:
            r.font.size = Pt(12.5); r.font.name = "Arial"
        p.space_after = Pt(9); p.line_spacing = 1.05

    # counter
    counter = DATA.get("counter", {})
    s = slide(); heading(s, "6 · Counter-evidence, risks & gaps", color=ALERT)
    tf = tb(s, .6, 1.0, 12.2, 6.3)
    for i, rk in enumerate(counter.get("risks", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = strip(rk.get("title", "")) + (" [WEAK SOURCE] " if rk.get("weak") else " ")
        r.font.bold = True; r.font.color.rgb = AMBER if rk.get("weak") else INK
        r2 = p.add_run(); r2.text = strip(rk.get("body", "")); r2.font.color.rgb = BODY
        for r in p.runs:
            r.font.size = Pt(10.5); r.font.name = "Arial"
        p.space_after = Pt(8); p.line_spacing = 1.02

    # sources
    srcs = DATA.get("sources", [])
    s = slide(); heading(s, "7 · Sources — all live-verified (%d)" % len(srcs))
    half = (len(srcs) + 1) // 2

    def fill_src(tf, items):
        first = True
        for x in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            r = p.add_run(); r.text = "%s. %s — %s" % (x.get("n"), strip(x.get("pub", "")), strip(x.get("title", "")))
            r.font.size = Pt(7.5); r.font.color.rgb = BODY; r.font.name = "Arial"; r.font.bold = True
            p.space_after = Pt(1)
            p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = x.get("url", "")
            rr.font.size = Pt(7); rr.font.color.rgb = GRN; rr.font.name = "Arial"
            p2.space_after = Pt(4)
    fill_src(tb(s, .5, 1.0, 6.1, 6.3), srcs[:half])
    fill_src(tb(s, 6.8, 1.0, 6.1, 6.3), srcs[half:])

    prs.save(str(out_path))


if __name__ == "__main__":
    main()
